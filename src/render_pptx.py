#!/usr/bin/env python3
"""storyboard.json → 실무 스타일 화면설계서 PPTX 렌더러 (v0.1)

표준: standards/style-spec.md (실무 문서 6종 기반)
사용: .venv/bin/python builder/render_pptx.py projects/<이름>/storyboard.json [-o out.pptx]
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── 팔레트 (style-spec §0, §5) ──────────────────────────────────────────
DARK = RGBColor(0x44, 0x44, 0x44)        # 헤더바/GNB 밴드
COVER_BG = RGBColor(0x58, 0x59, 0x5B)    # 표지 다크그레이
ORANGE = RGBColor(0xFF, 0x5A, 0x00)
MINT = RGBColor(0x17, 0xE5, 0xB2)
PURPLE = RGBColor(0x8E, 0x3F, 0xA8)
SKY = RGBColor(0x6F, 0xD8, 0xF2)
GRAY_DECO = RGBColor(0xA7, 0xA9, 0xAC)
RED = RGBColor(0xE5, 0x39, 0x35)         # 번호 마커
BADGE_BLUE = RGBColor(0x4A, 0x90, 0xD9)  # 수정사항 배지
LINK_BLUE = RGBColor(0x25, 0x63, 0xC0)
INK = RGBColor(0x22, 0x22, 0x22)
GRAY_TXT = RGBColor(0x77, 0x77, 0x77)
TBL_HEAD = RGBColor(0xE8, 0xE8, 0xE8)    # 연회색 표 헤더
TBL_LABEL = RGBColor(0xD9, 0xD9, 0xD9)   # 개요 라벨 열
WF_BG = RGBColor(0xF4, 0xF4, 0xF4)       # 와이어프레임 연회색 영역
WF_BOX = RGBColor(0xD9, 0xD9, 0xD9)
WF_INFO = RGBColor(0xE8, 0xEE, 0xF7)     # 연파랑 안내 박스
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0x9A, 0x9A, 0x9A)
HL_YELLOW = "FFF200"

FONT = "맑은 고딕"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── 저수준 헬퍼 ─────────────────────────────────────────────────────────

def _set_font(run, size=10, bold=False, color=INK, italic=False, underline=False, strike=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.underline = underline
    f.color.rgb = color
    ea = run._r.get_or_add_rPr().makeelement(qn("a:ea"), {"typeface": FONT})
    run._r.get_or_add_rPr().append(ea)
    if strike:
        run._r.get_or_add_rPr().set("strike", "sngStrike")


def _highlight(run, color=HL_YELLOW):
    rPr = run._r.get_or_add_rPr()
    hl = rPr.makeelement(qn("a:highlight"), {})
    clr = rPr.makeelement(qn("a:srgbClr"), {"val": color})
    hl.append(clr)
    rPr.append(hl)


_TOKEN = re.compile(r"(\*\*.+?\*\*|==.+?==|~~.+?~~|\[\[.+?\]\]|\{.+?\})")


def add_rich(par, text, size=10, color=INK, bold=False, italic=False):
    """마크업 지원 텍스트: **볼드** ==노랑== ~~취소~~ [[링크]] {변수=노랑}"""
    for tok in _TOKEN.split(text):
        if not tok:
            continue
        r = par.add_run()
        if tok.startswith("**") and tok.endswith("**"):
            r.text = tok[2:-2]
            _set_font(r, size, True, color, italic)
        elif tok.startswith("==") and tok.endswith("=="):
            r.text = tok[2:-2]
            _set_font(r, size, bold, color, italic)
            _highlight(r)
        elif tok.startswith("~~") and tok.endswith("~~"):
            r.text = tok[2:-2]
            _set_font(r, size, bold, GRAY_TXT, italic, strike=True)
        elif tok.startswith("[[") and tok.endswith("]]"):
            r.text = tok[2:-2]
            _set_font(r, size, bold, LINK_BLUE, italic, underline=True)
        elif tok.startswith("{") and tok.endswith("}"):
            r.text = tok
            _set_font(r, size, bold, color, italic)
            _highlight(r)
        else:
            r.text = tok
            _set_font(r, size, bold, color, italic)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, text, size=10, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(18000)
    tf.margin_top = tf.margin_bottom = Emu(9000)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        add_rich(p, ln, size=size, color=color, bold=bold, italic=italic)
    return tb


def dashed_line(sp):
    ln = sp.line._get_or_add_ln()
    d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
    ln.append(d)


# ── 공통 프레임 ─────────────────────────────────────────────────────────

def deco_triangles(slide, right=Inches(13.05), top=Inches(0.10)):
    """헤더 우측 지오메트릭 장식"""
    size = Inches(0.20)
    xs = [right - Inches(1.05), right - Inches(0.80), right - Inches(0.55), right - Inches(0.30)]
    colors = [ORANGE, MINT, GRAY_DECO, SKY]
    for i, (x, c) in enumerate(zip(xs, colors)):
        sp = rect(slide, x, top + (Inches(0.10) if i % 2 else 0), size, size,
                  fill=c, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
        if i % 2:
            sp.rotation = 180


def header_bar(slide, title):
    rect(slide, 0, 0, SLIDE_W, Inches(0.52), fill=DARK)
    txt(slide, Inches(0.25), Inches(0.02), Inches(9.5), Inches(0.48), title,
        size=17, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    deco_triangles(slide)


def page_number(slide, n):
    rect(slide, Inches(12.85), Inches(7.14), Inches(0.30), Inches(0.30),
         fill=RGBColor(0xDD, 0xDD, 0xDD), shape=MSO_SHAPE.PARALLELOGRAM)
    txt(slide, Inches(12.85), Inches(7.10), Inches(0.45), Inches(0.35), str(n),
        size=10, color=GRAY_TXT, italic=True, align=PP_ALIGN.CENTER)


def badge(slide, text, x, y):
    w = Inches(0.16 + 0.105 * len(text))
    rect(slide, x, y, w, Inches(0.26), fill=BADGE_BLUE)
    txt(slide, x, y - Inches(0.015), w, Inches(0.28), text, size=9, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── 페이지: 표지 ────────────────────────────────────────────────────────

def render_cover(prs, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COVER_BG)
    # 오렌지 대각 컷 (우측)
    tri = rect(slide, Inches(7.9), Inches(0), Inches(5.45), Inches(7.5),
               fill=ORANGE, shape=MSO_SHAPE.RIGHT_TRIANGLE)
    tri.rotation = 180
    # 지오메트릭 도형
    for (x, y, w, h, c, shp, rot) in [
        (9.3, 2.6, 0.55, 0.42, MINT, MSO_SHAPE.PARALLELOGRAM, 0),
        (12.2, 2.55, 0.55, 0.42, GRAY_DECO, MSO_SHAPE.PARALLELOGRAM, 0),
        (12.75, 3.3, 0.42, 0.42, PURPLE, MSO_SHAPE.PARALLELOGRAM, 0),
        (10.2, 3.7, 0.5, 0.45, MINT, MSO_SHAPE.ISOSCELES_TRIANGLE, 0),
        (8.4, 4.15, 0.55, 0.42, ORANGE, MSO_SHAPE.PARALLELOGRAM, 0),
        (11.2, 4.3, 0.5, 0.45, SKY, MSO_SHAPE.ISOSCELES_TRIANGLE, 0),
        (11.7, 4.35, 0.55, 0.42, GRAY_DECO, MSO_SHAPE.PARALLELOGRAM, 0),
        (12.3, 4.35, 0.42, 0.42, PURPLE, MSO_SHAPE.ISOSCELES_TRIANGLE, 180),
        (10.0, 4.85, 0.55, 0.45, PURPLE, MSO_SHAPE.PARALLELOGRAM, 0),
        (11.05, 4.9, 0.55, 0.42, MINT, MSO_SHAPE.PARALLELOGRAM, 0),
    ]:
        sp = rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill=c, shape=shp)
        sp.rotation = rot
    # 텍스트
    txt(slide, Inches(0.62), Inches(2.15), Inches(4), Inches(0.4),
        meta.get("company", "Smilegate") + "˚", size=17, color=RGBColor(0xB9, 0xBB, 0xBE), bold=True)
    txt(slide, Inches(0.62), Inches(2.62), Inches(6), Inches(0.45),
        meta.get("category", ""), size=20, color=ORANGE, bold=True)
    txt(slide, Inches(0.60), Inches(3.05), Inches(8.2), Inches(0.9),
        meta["title"], size=32, color=WHITE, bold=True)
    rows = [("작성자", meta.get("author", "")), ("소속", meta.get("team", "")),
            ("최초 작성일", meta.get("created", "")), ("최종 업데이트", meta.get("updated", ""))]
    y = 4.35
    for label, val in rows:
        txt(slide, Inches(0.66), Inches(y), Inches(1.3), Inches(0.3), label, size=10.5, color=WHITE)
        txt(slide, Inches(1.85), Inches(y), Inches(3.5), Inches(0.3), val, size=10.5, color=WHITE, bold=True)
        y += 0.295
    txt(slide, Inches(0.62), Inches(6.95), Inches(6), Inches(0.3),
        meta.get("copyright", "©Smilegate Megaport Inc. All Rights Reserved."),
        size=9, color=RGBColor(0x9d, 0x9f, 0xa2))
    return slide


# ── 페이지: History ────────────────────────────────────────────────────

def render_history(prs, history, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "History")
    n = len(history)
    tbl_shape = slide.shapes.add_table(n + 1, 3, Inches(0.38), Inches(0.72),
                                       Inches(12.55), Inches(0.34 * (n + 1)))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.2)
    tbl.columns[1].width = Inches(10.15)
    tbl.columns[2].width = Inches(1.2)
    _strip_table_style(tbl_shape)
    heads = ["Date", "History", "Drafter"]
    for c, head in enumerate(heads):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = TBL_HEAD
        _cell_text(cell, [head], size=10.5, bold=False, align=PP_ALIGN.CENTER)
    for r, h in enumerate(history, start=1):
        c0 = tbl.cell(r, 0); c0.fill.solid(); c0.fill.fore_color.rgb = WHITE
        _cell_text(c0, [h["date"]], size=9.5, align=PP_ALIGN.CENTER)
        c1 = tbl.cell(r, 1); c1.fill.solid(); c1.fill.fore_color.rgb = WHITE
        _cell_text(c1, ["•  " + it for it in h["items"]], size=9)
        c2 = tbl.cell(r, 2); c2.fill.solid(); c2.fill.fore_color.rgb = WHITE
        _cell_text(c2, [h["drafter"]], size=9.5, align=PP_ALIGN.CENTER)
    page_number(slide, page_no)
    return slide


def _strip_table_style(graphic_frame):
    """기본 파란 표 스타일 제거 → 검정 얇은 테두리"""
    tbl = graphic_frame._element.graphic.graphicData.tbl
    tblPr = tbl.tblPr
    tblPr.set("firstRow", "0"); tblPr.set("bandRow", "0")
    for el in list(tblPr):
        tblPr.remove(el)


def _cell_text(cell, lines, size=10, bold=False, align=PP_ALIGN.LEFT, color=INK):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(45720); tf.margin_right = Emu(45720)
    tf.margin_top = Emu(27432); tf.margin_bottom = Emu(27432)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        add_rich(p, ln, size=size, color=color, bold=bold)


def _cell_border(cell, color="000000", w=6350):
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = tcPr.makeelement(qn(tag), {"w": str(w), "cap": "flat"})
        fill = tcPr.makeelement(qn("a:solidFill"), {})
        clr = tcPr.makeelement(qn("a:srgbClr"), {"val": color})
        fill.append(clr); ln.append(fill); tcPr.append(ln)


def _all_borders(tbl):
    for row in tbl.rows:
        for cell in row.cells:
            _cell_border(cell)


# ── 페이지: 개요 ────────────────────────────────────────────────────────

def render_overview(prs, overview, page_no, title="개요"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, title)
    n = len(overview)
    total_h = 6.1
    tbl_shape = slide.shapes.add_table(n, 2, Inches(0.38), Inches(0.85),
                                       Inches(12.55), Inches(total_h))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.7)
    tbl.columns[1].width = Inches(10.85)
    _strip_table_style(tbl_shape)
    weights = [max(1, len(r["lines"])) for r in overview]
    tw = sum(weights)
    for r, rowdef in enumerate(overview):
        tbl.rows[r].height = Inches(total_h * weights[r] / tw)
        c0 = tbl.cell(r, 0); c0.fill.solid(); c0.fill.fore_color.rgb = TBL_LABEL
        _cell_text(c0, [rowdef["label"]], size=11, bold=True, align=PP_ALIGN.CENTER)
        c1 = tbl.cell(r, 1); c1.fill.solid(); c1.fill.fore_color.rgb = WHITE
        _cell_text(c1, [("•  " + ln if not ln.startswith(" ") else "     •  " + ln.strip())
                        for ln in rowdef["lines"]], size=10)
    _all_borders(tbl)
    page_number(slide, page_no)
    return slide


# ── 와이어프레임 요소 렌더 ──────────────────────────────────────────────

def _wf_coord(cx, cy, cw, ch, e):
    """캔버스 상대좌표(0~100) → EMU"""
    return (Emu(int(cx + cw * e["x"] / 100)), Emu(int(cy + ch * e["y"] / 100)),
            Emu(int(cw * e["w"] / 100)), Emu(int(ch * e["h"] / 100)))


def render_wf_element(slide, cx, cy, cw, ch, e):
    x, y, w, h = _wf_coord(cx, cy, cw, ch, e)
    kind = e["kind"]
    text = e.get("text", "")
    variant = e.get("variant", "")
    size = e.get("size", 9)
    align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(
        e.get("align", "center"), PP_ALIGN.CENTER)

    if kind == "band":
        rect(slide, x, y, w, h, fill=DARK)
        if text:
            txt(slide, x, y, w, h, text, size=size + 2, color=RGBColor(0x99, 0x99, 0x99),
                bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "box":
        fill = {"light": WF_BG, "info": WF_INFO, "white": WHITE, "gray": WF_BOX}.get(variant, WF_BG)
        line = None if variant == "noline" else BORDER
        rect(slide, x, y, w, h, fill=fill, line=line)
        if text:
            txt(slide, x, y, w, h, text, size=size, align=align, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "image_ph":
        rect(slide, x, y, w, h, fill=WF_BOX, line=BORDER)
        txt(slide, x, y, w, h, text or "아이콘\nIMG", size=size, color=RGBColor(0x66, 0x66, 0x66),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "text":
        txt(slide, x, y, w, h, text, size=size, align=align,
            bold=(variant == "bold"), anchor=MSO_ANCHOR.MIDDLE,
            color=LINK_BLUE if variant == "link" else INK)
    elif kind == "button":
        if variant == "outline":
            rect(slide, x, y, w, h, fill=WHITE, line=BORDER)
            txt(slide, x, y, w, h, text, size=size, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            sp = rect(slide, x, y, w, h, fill=DARK)
            sp.shadow.inherit = True
            txt(slide, x, y, w, h, text, size=size, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "tabs":
        cols = e.get("columns", [])
        active = int(e.get("variant") or 0)
        tw_ = int(w / max(1, len(cols)))
        for i, cname in enumerate(cols):
            fill = DARK if i == active else WF_BOX
            color = WHITE if i == active else RGBColor(0x55, 0x55, 0x55)
            rect(slide, Emu(int(x) + tw_ * i), y, Emu(tw_ - 12700), h, fill=fill, line=BORDER)
            txt(slide, Emu(int(x) + tw_ * i), y, Emu(tw_ - 12700), h, cname, size=size,
                color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "chips":
        cols = e.get("columns", [])
        active = e.get("variant")
        active = int(active) if active not in (None, "") else -1
        gap = 27432
        cwidth = int((int(w) - gap * (len(cols) - 1)) / max(1, len(cols)))
        for i, cname in enumerate(cols):
            fill = DARK if i == active else WHITE
            color = WHITE if i == active else INK
            xx = Emu(int(x) + (cwidth + gap) * i)
            rect(slide, xx, y, Emu(cwidth), h, fill=fill, line=BORDER)
            txt(slide, xx, y, Emu(cwidth), h, cname, size=size - 1, color=color,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "table":
        cols = e.get("columns", [])
        rows = e.get("rows", [])
        nrow = 1 + len(rows)
        rh = int(h / nrow)
        colw = int(w / max(1, len(cols)))
        for c, cname in enumerate(cols):
            xx = Emu(int(x) + colw * c)
            rect(slide, xx, y, Emu(colw), Emu(rh), fill=DARK, line=WHITE)
            txt(slide, xx, y, Emu(colw), Emu(rh), cname, size=size - 1, color=WHITE,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        for r, rowvals in enumerate(rows):
            yy = Emu(int(y) + rh * (r + 1))
            for c in range(len(cols)):
                val = rowvals[c] if c < len(rowvals) else ""
                xx = Emu(int(x) + colw * c)
                rect(slide, xx, yy, Emu(colw), Emu(rh), fill=WHITE, line=BORDER)
                txt(slide, xx, yy, Emu(colw), Emu(rh), val, size=size - 1,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "amount_bar":
        rect(slide, x, y, w, h, fill=RGBColor(0xE2, 0xE2, 0xE2))
        txt(slide, x, y, Emu(int(w) - 91440), h, text, size=size + 5, bold=True,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "input":
        rect(slide, x, y, w, h, fill=WHITE, line=BORDER)
        txt(slide, Emu(int(x) + 45720), y, w, h, text, size=size, color=GRAY_TXT,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "note":
        txt(slide, x, y, w, h, text, size=size, color=INK, align=align, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "dashed":
        sp = rect(slide, x, y, w, h, fill=None, line=RED if variant == "red" else BADGE_BLUE, line_w=1.2)
        dashed_line(sp)
    elif kind == "diamond":
        rect(slide, x, y, w, h, fill=WHITE, line=INK, shape=MSO_SHAPE.DIAMOND)
        txt(slide, x, y, w, h, text, size=size - 1, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "rounded":
        fill = DARK if variant == "dark" else WHITE
        color = WHITE if variant == "dark" else INK
        rect(slide, x, y, w, h, fill=fill, line=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(slide, x, y, w, h, text, size=size, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "db":
        rect(slide, x, y, w, h, fill=WF_BOX, line=INK, shape=MSO_SHAPE.CAN)
        txt(slide, x, Emu(int(y) + int(h) // 5), w, h, text, size=size - 1,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    elif kind == "dim":
        sp = rect(slide, x, y, w, h, fill=RGBColor(0x11, 0x11, 0x11))
        # 55% 투명
        fill_el = sp.fill.fore_color._xFill
        srgb = fill_el.find(qn("a:srgbClr"))
        alpha = srgb.makeelement(qn("a:alpha"), {"val": "45000"})
        srgb.append(alpha)
    elif kind == "arrow":
        render_wf_arrow(slide, cx, cy, cw, ch, e)


def render_wf_arrow(slide, cx, cy, cw, ch, e):
    """kind=arrow 전용: x,y(시작) → x2,y2(끝), label"""
    from pptx.enum.shapes import MSO_CONNECTOR
    x1 = Emu(int(cx + cw * e["x"] / 100)); y1 = Emu(int(cy + ch * e["y"] / 100))
    x2 = Emu(int(cx + cw * e["x2"] / 100)); y2 = Emu(int(cy + ch * e["y2"] / 100))
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = RED if e.get("variant") == "red" else INK
    conn.line.width = Pt(1.1)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "arrow"}))
    conn.shadow.inherit = False
    if e.get("text"):
        mx = Emu((int(x1) + int(x2)) // 2 - 137160)
        my = Emu((int(y1) + int(y2)) // 2 - 91440)
        txt(slide, mx, my, Inches(0.45), Inches(0.22), e["text"], size=8,
            bold=True, align=PP_ALIGN.CENTER)


def render_marker(slide, cx, cy, cw, ch, m):
    mx = Emu(int(cx + cw * m["x"] / 100))
    my = Emu(int(cy + ch * m["y"] / 100))
    s = Inches(0.21)
    label = str(m["no"])
    w = s if len(label) <= 1 else Inches(0.30)
    rect(slide, mx, my, w, s, fill=RED)
    txt(slide, mx, my - Inches(0.02), w, Inches(0.24), label, size=10, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, italic=True)


# ── 페이지: PC 화면 상세 ───────────────────────────────────────────────

WF_X = Inches(0.35)
WF_Y = Inches(0.78)
WF_W = Inches(8.5)
WF_H = Inches(6.35)
DESC_X = Inches(10.30)
DESC_W = Inches(2.70)


def render_screen_pc(prs, page, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, page.get("section", ""))

    wf = page.get("wireframe", {})
    cx, cy, cw, ch = int(WF_X), int(WF_Y), int(WF_W), int(WF_H)
    # 캔버스 외곽
    rect(slide, WF_X, WF_Y, WF_W, WF_H, fill=WHITE, line=DARK, line_w=1.1)
    for e in wf.get("elements", []):
        render_wf_element(slide, cx, cy, cw, ch, e)
    for m in wf.get("markers", []):
        render_marker(slide, cx, cy, cw, ch, m)
    if wf.get("omitted_below"):
        yy = Emu(cy + ch + 27432)
        sp = rect(slide, WF_X, yy, WF_W, Emu(1), fill=None, line=GRAY_TXT, line_w=0.75)
        dashed_line(sp)
        txt(slide, Inches(3.9), Emu(int(yy) + 18000), Inches(1.5), Inches(0.25), "이하 생략",
            size=8.5, color=GRAY_TXT, italic=True, align=PP_ALIGN.CENTER)

    # Description 테이블
    descs = page.get("descriptions", [])
    if descs:
        n = len(descs)
        tbl_shape = slide.shapes.add_table(n + 1, 2, DESC_X, Inches(0.72),
                                           DESC_W, Inches(0.3))
        tbl = tbl_shape.table
        tbl.columns[0].width = Inches(0.38)
        tbl.columns[1].width = Inches(2.32)
        _strip_table_style(tbl_shape)
        h0 = tbl.cell(0, 0); h0.fill.solid(); h0.fill.fore_color.rgb = DARK
        _cell_text(h0, ["No"], size=9, color=WHITE, align=PP_ALIGN.CENTER)
        h1 = tbl.cell(0, 1); h1.fill.solid(); h1.fill.fore_color.rgb = DARK
        _cell_text(h1, ["Description"], size=9, color=WHITE, align=PP_ALIGN.CENTER)
        for r, d in enumerate(descs, start=1):
            c0 = tbl.cell(r, 0); c0.fill.solid(); c0.fill.fore_color.rgb = WHITE
            _cell_text(c0, [str(d["no"])], size=9, align=PP_ALIGN.CENTER, color=GRAY_TXT)
            c1 = tbl.cell(r, 1); c1.fill.solid(); c1.fill.fore_color.rgb = WHITE
            lines = []
            title = "**" + d["title"] + "**"
            if d.get("etype"):
                title += " **[" + d["etype"] + "]**"
            lines.append(title)
            for ln in d.get("lines", []):
                if ln.startswith("> "):
                    lines.append("      ➢  " + ln[2:])
                else:
                    lines.append(ln)
            _cell_text(c1, lines, size=8)
        _all_borders(tbl)

    if page.get("badge"):
        badge(slide, page["badge"], DESC_X + Inches(0.9), Inches(0.72) + Inches(0.3) * (len(descs) + 1) + Inches(0.35))
    page_number(slide, page_no)
    return slide


# ── 페이지: 정책 ────────────────────────────────────────────────────────

def _render_block_table(slide, x, y, w, block):
    cols = block["columns"]
    rows = block.get("rows", [])
    variant = block.get("variant", "gray")
    col_ws = block.get("col_widths")  # 비율 리스트
    row_h = Inches(block.get("row_h", 0.30))
    n = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n, len(cols), x, y, w, Emu(int(row_h) * n))
    tbl = tbl_shape.table
    if col_ws:
        total = sum(col_ws)
        for i, cwr in enumerate(col_ws):
            tbl.columns[i].width = Emu(int(int(w) * cwr / total))
    _strip_table_style(tbl_shape)
    head_fill = DARK if variant == "dark" else TBL_HEAD
    head_color = WHITE if variant == "dark" else INK
    for c, cname in enumerate(cols):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = head_fill
        _cell_text(cell, [cname], size=9, color=head_color, align=PP_ALIGN.CENTER, bold=(variant != "dark"))
    for r, rowvals in enumerate(rows, start=1):
        for c in range(len(cols)):
            cell = tbl.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            _cell_text(cell, [rowvals[c] if c < len(rowvals) else ""], size=9)
    _all_borders(tbl)
    return Emu(int(y) + int(row_h) * n)


def render_policy(prs, page, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, page.get("section", ""))
    if page.get("badge"):
        badge(slide, page["badge"], Inches(11.55), Inches(0.13))
    y = Inches(0.75)
    for block in page.get("blocks", []):
        kind = block["kind"]
        if kind == "bullets":
            items = block["items"]
            box_h = Inches(0.24 * len(items) + 0.1)
            tb = slide.shapes.add_textbox(Inches(0.35), y, Inches(12.6), box_h)
            tf = tb.text_frame
            tf.word_wrap = True
            for i, it in enumerate(items):
                text = it["text"] if isinstance(it, dict) else it
                level = it.get("level", 0) if isinstance(it, dict) else 0
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                indent = "        " * level
                add_rich(p, indent + "•   " + text, size=10.5 - level, color=INK)
            y = Emu(int(y) + int(box_h) + 91440)
        elif kind == "table":
            tx = Inches(block.get("x", 0.5))
            tw = Inches(block.get("w", 6.5))
            y_end = _render_block_table(slide, tx, y, tw, block)
            if not block.get("float"):
                y = Emu(int(y_end) + 137160)
        elif kind == "spacer":
            y = Emu(int(y) + Inches(block.get("h", 0.2)))
    page_number(slide, page_no)
    return slide


# ── 페이지: 목차(섹션 구분자) ──────────────────────────────────────────

def render_section_index(prs, page, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, page.get("section", "Index"))
    items = page.get("items", [])
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(10), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        text = it["text"] if isinstance(it, dict) else it
        current = it.get("current", False) if isinstance(it, dict) else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = ("■  " if current else "□  ") + text
        _set_font(r, 16 if current else 14, bold=current,
                  color=INK if current else RGBColor(0xAA, 0xAA, 0xAA))
    page_number(slide, page_no)
    return slide


# ── 페이지: MO 화면 상세 ───────────────────────────────────────────────

MO_FRAME_W = Inches(2.35)
MO_FRAME_H = Inches(5.9)


def render_screen_mo(prs, page, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, page.get("section", ""))
    frames = page.get("wireframe", {}).get("frames", [])
    x0 = Inches(0.35)
    gap = Inches(0.35)
    for i, fr in enumerate(frames):
        fx = Emu(int(x0) + (int(MO_FRAME_W) + int(gap)) * i)
        fy = Inches(0.95)
        rect(slide, fx, fy, MO_FRAME_W, MO_FRAME_H, fill=WHITE, line=DARK, line_w=1.2)
        cx, cy, cw, ch = int(fx), int(fy), int(MO_FRAME_W), int(MO_FRAME_H)
        for e in fr.get("elements", []):
            render_wf_element(slide, cx, cy, cw, ch, e)
        for m in fr.get("markers", []):
            render_marker(slide, cx, cy, cw, ch, m)
        if fr.get("label"):
            txt(slide, fx, Inches(0.68), MO_FRAME_W, Inches(0.25), fr["label"],
                size=9, color=GRAY_TXT, align=PP_ALIGN.CENTER)
        if i < len(frames) - 1 and fr.get("continues", True):
            txt(slide, Emu(int(fx) + int(MO_FRAME_W) - 27432), Inches(6.9),
                Inches(0.9), Inches(0.25), "이어서 ▶", size=8, color=GRAY_TXT, italic=True)
    # MO는 Description 생략이 기본 — 있으면 우측에
    descs = page.get("descriptions", [])
    if descs:
        n = len(descs)
        tbl_shape = slide.shapes.add_table(n + 1, 2, DESC_X, Inches(0.72), DESC_W, Inches(0.3))
        tbl = tbl_shape.table
        tbl.columns[0].width = Inches(0.38)
        tbl.columns[1].width = Inches(2.32)
        _strip_table_style(tbl_shape)
        for c, head in enumerate(["No", "Description"]):
            cell = tbl.cell(0, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = DARK
            _cell_text(cell, [head], size=9, color=WHITE, align=PP_ALIGN.CENTER)
        for r, d in enumerate(descs, start=1):
            c0 = tbl.cell(r, 0); c0.fill.solid(); c0.fill.fore_color.rgb = WHITE
            _cell_text(c0, [str(d["no"])], size=9, align=PP_ALIGN.CENTER, color=GRAY_TXT)
            c1 = tbl.cell(r, 1); c1.fill.solid(); c1.fill.fore_color.rgb = WHITE
            lines = ["**" + d["title"] + ("** **[" + d["etype"] + "]**" if d.get("etype") else "**")]
            lines += [("      ➢  " + ln[2:]) if ln.startswith("> ") else ln for ln in d.get("lines", [])]
            _cell_text(c1, lines, size=8)
        _all_borders(tbl)
    else:
        txt(slide, Inches(10.4), Inches(0.95), Inches(2.6), Inches(0.6),
            "* Description은 PC 페이지 정의와 동일", size=9, color=GRAY_TXT, italic=True)
    page_number(slide, page_no)
    return slide


# ── 페이지: 의사결정 (A안/B안) ─────────────────────────────────────────

def render_decision(prs, page, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, page.get("section", ""))
    panels = page.get("panels", [])
    pw = Inches(6.0)
    ph = Inches(4.6)
    xs = [Inches(0.35), Inches(6.75)]
    for i, panel in enumerate(panels[:2]):
        x = xs[i]
        y = Inches(0.85)
        chosen = panel.get("chosen", False)
        border_color = RED if chosen else BORDER
        rect(slide, x, y, pw, ph, fill=WF_BG if not chosen else WHITE,
             line=border_color, line_w=2.5 if chosen else 1.0)
        txt(slide, Emu(int(x) + 91440), Emu(int(y) + 54864), Emu(int(pw) - 182880), Inches(0.35),
            panel.get("label", ""), size=11, bold=True)
        body_y = Emu(int(y) + 457200)
        tb = slide.shapes.add_textbox(Emu(int(x) + 137160), body_y,
                                      Emu(int(pw) - 274320), Emu(int(ph) - 548640))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, ln in enumerate(panel.get("lines", [])):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            add_rich(p, ln, size=10)
        if chosen and panel.get("stamp"):
            badge(slide, panel["stamp"], Emu(int(x) + int(pw) - Inches(1.7)), Emu(int(y) + 27432))
    ry = Inches(5.75)
    for ln in page.get("rationale", []):
        tb = slide.shapes.add_textbox(Inches(0.5), ry, Inches(12.3), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = ln
        _set_font(r, 10.5, bold=True, color=INK, underline=True)
        ry = Emu(int(ry) + 274320)
    if page.get("link_note"):
        tb = slide.shapes.add_textbox(Inches(0.5), ry, Inches(12.3), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        add_rich(p, "[[" + page["link_note"] + "]]", size=10.5, bold=True)
    page_number(slide, page_no)
    return slide


# ── 페이지: 어드민 (통합운영툴) ────────────────────────────────────────

ADMIN_BLUE = RGBColor(0x2F, 0x5B, 0x9E)


def render_admin(prs, page, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, page.get("section", ""))
    adm = page.get("admin", {})
    # 셸 캔버스
    AX, AY, AW, AH = WF_X, WF_Y, WF_W, WF_H
    rect(slide, AX, AY, AW, AH, fill=WHITE, line=DARK, line_w=1.1)
    # 파란 상단바
    rect(slide, AX, AY, AW, Inches(0.34), fill=ADMIN_BLUE)
    txt(slide, Emu(int(AX) + 91440), AY, Inches(4), Inches(0.34),
        adm.get("shell_title", "통합운영툴"), size=10, color=WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE)
    # 좌측 트리
    tree_w = Inches(1.55)
    ty = Emu(int(AY) + int(Inches(0.34)))
    rect(slide, AX, ty, tree_w, Emu(int(AH) - int(Inches(0.34))),
         fill=RGBColor(0xF7, 0xF7, 0xF7), line=BORDER)
    iy = Emu(int(ty) + 54864)
    for item in adm.get("tree", []):
        text = item["text"] if isinstance(item, dict) else item
        level = item.get("level", 0) if isinstance(item, dict) else 0
        current = item.get("current", False) if isinstance(item, dict) else False
        new = item.get("new", False) if isinstance(item, dict) else False
        ih = Inches(0.24)
        if current:
            rect(slide, Emu(int(AX) + 18000), iy, Emu(int(tree_w) - 36000), ih,
                 fill=RGBColor(0xDC, 0xE6, 0xF4))
        txt(slide, Emu(int(AX) + 45720 + level * 137160), Emu(int(iy) - 13716),
            Emu(int(tree_w) - 91440), Inches(0.26), text, size=8,
            bold=current, color=ADMIN_BLUE if current else INK)
        if new:
            sp = rect(slide, Emu(int(AX) + 27432), iy, Emu(int(tree_w) - 54864), ih,
                      fill=None, line=RED, line_w=1.0)
            dashed_line(sp)
        iy = Emu(int(iy) + int(ih) + 18000)
    # 콘텐츠 영역 = 트리 우측
    cx = int(AX) + int(tree_w)
    cy = int(ty)
    cw = int(AW) - int(tree_w)
    ch = int(AH) - int(Inches(0.34))
    for e in adm.get("elements", []):
        render_wf_element(slide, cx, cy, cw, ch, e)
    for m in adm.get("markers", []):
        render_marker(slide, cx, cy, cw, ch, m)
    # Description
    _desc_table(slide, page.get("descriptions", []))
    if page.get("badge"):
        badge(slide, page["badge"], Inches(11.55), Inches(0.13))
    page_number(slide, page_no)
    return slide


def _desc_table(slide, descs):
    if not descs:
        return
    n = len(descs)
    tbl_shape = slide.shapes.add_table(n + 1, 2, DESC_X, Inches(0.72), DESC_W, Inches(0.3))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(0.38)
    tbl.columns[1].width = Inches(2.32)
    _strip_table_style(tbl_shape)
    for c, head in enumerate(["No", "Description"]):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK
        _cell_text(cell, [head], size=9, color=WHITE, align=PP_ALIGN.CENTER)
    for r, d in enumerate(descs, start=1):
        c0 = tbl.cell(r, 0); c0.fill.solid(); c0.fill.fore_color.rgb = WHITE
        _cell_text(c0, [str(d["no"])], size=9, align=PP_ALIGN.CENTER, color=GRAY_TXT)
        c1 = tbl.cell(r, 1); c1.fill.solid(); c1.fill.fore_color.rgb = WHITE
        lines = ["**" + d["title"] + ("** **[" + d["etype"] + "]**" if d.get("etype") else "**")]
        lines += [("      ➢  " + ln[2:]) if ln.startswith("> ") else ln for ln in d.get("lines", [])]
        _cell_text(c1, lines, size=8)
    _all_borders(tbl)


# ── 페이지: 플로우차트 ─────────────────────────────────────────────────

def render_flowchart(prs, page, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, page.get("section", ""))
    # 좌측 순서도 캔버스 (외곽선 없음)
    FX, FY = Inches(0.35), Inches(0.78)
    FW, FH = Inches(7.6), Inches(6.4)
    cx, cy, cw, ch = int(FX), int(FY), int(FW), int(FH)
    wf = page.get("wireframe", {})
    for e in wf.get("elements", []):
        render_wf_element(slide, cx, cy, cw, ch, e)
    for m in wf.get("markers", []):
        render_marker(slide, cx, cy, cw, ch, m)
    # 우측: 얼럿 정책 표
    ax = Inches(8.25)
    y = Inches(0.85)
    at = page.get("alert_table")
    if at:
        txt(slide, ax, y, Inches(3), Inches(0.28), at.get("title", "■ 얼럿 정책"),
            size=10.5, bold=True)
        y = Emu(int(y) + 274320)
        y = _render_block_table(slide, ax, y, Inches(4.7),
                                {"columns": at["columns"], "rows": at["rows"],
                                 "variant": "dark", "row_h": at.get("row_h", 0.42),
                                 "col_widths": at.get("col_widths")})
        y = Emu(int(y) + 182880)
    # 우측: 팝업 템플릿 목업
    pt = page.get("popup_template")
    if pt:
        txt(slide, ax, y, Inches(3), Inches(0.28), "■ 레이어 팝업 템플릿", size=10.5, bold=True)
        y = Emu(int(y) + 274320)
        pw, ph = Inches(2.9), Inches(1.7)
        rect(slide, ax, y, pw, ph, fill=WHITE, line=INK, line_w=1.2)
        txt(slide, ax, Emu(int(y) + 137160), pw, Inches(0.8), pt.get("body", "{얼럿 메시지}"),
            size=9, align=PP_ALIGN.CENTER)
        bw = Inches(0.9)
        rect(slide, Emu(int(ax) + (int(pw) - int(bw)) // 2), Emu(int(y) + int(ph) - 411480),
             bw, Inches(0.3), fill=DARK)
        txt(slide, Emu(int(ax) + (int(pw) - int(bw)) // 2), Emu(int(y) + int(ph) - 429480),
            bw, Inches(0.32), pt.get("button", "{버튼명}"), size=8.5, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _desc_table(slide, page.get("descriptions", []))
    page_number(slide, page_no)
    return slide


# ── 메인 ────────────────────────────────────────────────────────────────

RENDERERS = {
    "screen_pc": render_screen_pc,
    "screen_mo": render_screen_mo,
    "policy": render_policy,
    "section_index": render_section_index,
    "decision": render_decision,
    "admin": render_admin,
    "flowchart": render_flowchart,
}


def build(data, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    render_cover(prs, data["meta"])
    page_no = 2
    if data.get("history"):
        render_history(prs, data["history"], page_no)
        page_no += 1
    if data.get("overview"):
        render_overview(prs, data["overview"], page_no)
        page_no += 1
    for page in data.get("pages", []):
        fn = RENDERERS.get(page["type"])
        if fn is None:
            print(f"⚠️  미구현 페이지 유형 스킵: {page['type']}")
            continue
        fn(prs, page, page_no)
        page_no += 1

    prs.save(out_path)
    print(f"✅ {out_path} ({page_no - 1} slides)")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("json_path")
    ap.add_argument("-o", "--out")
    args = ap.parse_args()
    src = Path(args.json_path)
    data = json.loads(src.read_text(encoding="utf-8"))
    ver = data.get("meta", {}).get("version", "0.1")
    out = args.out or str(src.parent / f"{data['meta']['title']}_v{ver}.pptx")
    build(data, out)
